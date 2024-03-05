from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Add a title slide
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]
title.text = "Data Science Capstone Project"
subtitle.text = "An Overview of the Project"

# Add content slides
content_slides = [
    ("Data Exploration", "Summarize your data exploration process."),
    ("Data Preprocessing", "Discuss the preprocessing steps you performed."),
    ("Model Selection", "Describe the machine learning models you experimented with."),
    # Add more slides as needed
]

for slide_title, slide_content in content_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = slide_title
    content.text = slide_content

# Save the presentation
prs.save("data_science_capstone.pptx")
